"""Headless read/write for Office documents (Excel / Word / PowerPoint).

A top automation need is reading and writing spreadsheets and documents
without driving the GUI. This module wraps the de-facto libraries
(``openpyxl`` / ``python-docx`` / ``python-pptx``) behind a small,
serialisable API so flows can ingest an ``.xlsx`` row-set or emit a
``.docx`` report headlessly.

Those libraries are an **optional** dependency: install them with
``pip install je_auto_control[office]``. Each function raises a clear
:class:`RuntimeError` when the backing library is missing, so the core
package stays lean and import-time stays Qt-free / dependency-free.
"""
from pathlib import Path
from typing import Any, Dict, List

_HINT = "pip install je_auto_control[office]"


def _openpyxl() -> Any:
    """Import openpyxl (optional Excel backend) or raise a helpful error."""
    try:
        import openpyxl
    except ImportError as error:
        raise RuntimeError(f"Excel I/O requires openpyxl ({_HINT}).") from error
    return openpyxl


def _docx() -> Any:
    """Import python-docx (optional Word backend) or raise a helpful error."""
    try:
        import docx
    except ImportError as error:
        raise RuntimeError(
            f"Word I/O requires python-docx ({_HINT}).") from error
    return docx


def _pptx() -> Any:
    """Import python-pptx (optional PPT backend) or raise a helpful error."""
    try:
        import pptx
    except ImportError as error:
        raise RuntimeError(
            f"PowerPoint I/O requires python-pptx ({_HINT}).") from error
    return pptx


def _existing(path: str) -> Path:
    resolved = Path(path).expanduser()
    if not resolved.is_file():
        raise FileNotFoundError(f"no such file: {resolved}")
    return resolved


# --- Excel (.xlsx) --------------------------------------------------------

def read_workbook(path: str, sheet: str = "") -> List[Dict[str, Any]]:
    """Read a worksheet into a list of dicts (first row supplies the keys).

    ``sheet`` defaults to the active sheet.
    """
    openpyxl = _openpyxl()
    workbook = openpyxl.load_workbook(filename=str(_existing(path)),
                                      read_only=True, data_only=True)
    try:
        worksheet = workbook[sheet] if sheet else workbook.active
        rows_iter = worksheet.iter_rows(values_only=True)
        header = next(rows_iter, None)
        if header is None:
            return []
        keys = [str(cell) for cell in header]
        return [dict(zip(keys, values)) for values in rows_iter]
    finally:
        workbook.close()


def write_workbook(path: str, rows: List[Dict[str, Any]],
                   sheet: str = "Sheet1") -> str:
    """Write ``rows`` (list of dicts) to an ``.xlsx`` file; return the path."""
    openpyxl = _openpyxl()
    workbook = openpyxl.Workbook()
    worksheet = workbook.active
    worksheet.title = sheet
    rows = list(rows)
    if rows:
        keys = list(rows[0].keys())
        worksheet.append(keys)
        for row in rows:
            worksheet.append([row.get(key) for key in keys])
    target = Path(path).expanduser()
    workbook.save(str(target))
    return str(target.resolve())


# --- Word (.docx) ---------------------------------------------------------

def read_document(path: str) -> Dict[str, List[str]]:
    """Read a ``.docx`` file's paragraph texts."""
    docx = _docx()
    document = docx.Document(str(_existing(path)))
    return {"paragraphs": [para.text for para in document.paragraphs]}


def write_document(path: str, paragraphs: List[str]) -> str:
    """Write ``paragraphs`` to a ``.docx`` file; return the path."""
    docx = _docx()
    document = docx.Document()
    for paragraph in paragraphs:
        document.add_paragraph(str(paragraph))
    target = Path(path).expanduser()
    document.save(str(target))
    return str(target.resolve())


# --- PowerPoint (.pptx) ---------------------------------------------------

def read_presentation(path: str) -> Dict[str, List[List[str]]]:
    """Read a ``.pptx`` file's per-slide text runs."""
    pptx = _pptx()
    presentation = pptx.Presentation(str(_existing(path)))
    slides = []
    for slide in presentation.slides:
        slides.append([shape.text for shape in slide.shapes
                       if shape.has_text_frame and shape.text])
    return {"slides": slides}


def _add_slide(presentation: Any, layout: Any, spec: Any) -> None:
    slide = presentation.slides.add_slide(layout)
    title = spec.get("title", "") if isinstance(spec, dict) else str(spec)
    body = spec.get("body", []) if isinstance(spec, dict) else []
    if slide.shapes.title is not None:
        slide.shapes.title.text = str(title)
    if body:
        frame = slide.placeholders[1].text_frame
        frame.text = str(body[0])
        for line in body[1:]:
            frame.add_paragraph().text = str(line)


def write_presentation(path: str, slides: List[Any]) -> str:
    """Write ``slides`` (each ``{title, body:[...]}``) to a ``.pptx`` file."""
    pptx = _pptx()
    presentation = pptx.Presentation()
    layout = presentation.slide_layouts[1]   # "Title and Content"
    for spec in slides:
        _add_slide(presentation, layout, spec)
    target = Path(path).expanduser()
    presentation.save(str(target))
    return str(target.resolve())
